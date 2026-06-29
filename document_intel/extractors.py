# document_intel/extractors.py
# ─────────────────────────────────────────────────────────────────────────
# Extracts text/tables from uploaded business documents (PDF, DOCX, PPTX,
# XLSX, CSV, TXT) into a single traceable structure:
#
#   ExtractedDocument(
#       filename, file_type,
#       units = [
#           ExtractedUnit(locator="Page 7",  text="...", kind="text"),
#           ExtractedUnit(locator="Page 7",  text="| A | B |...", kind="table"),
#           ExtractedUnit(locator="Slide 3", text="...", kind="slide"),
#           ExtractedUnit(locator="Sheet: Revenue", text="...", kind="sheet"),
#           ...
#       ]
#   )
#
# Every unit carries a `locator` string (page/slide/sheet reference) so the
# contradiction detector and the unified-context builder can always cite
# "Source File, Page 7" / "Slide 3" / "Sheet: Revenue" back to the user.
# ─────────────────────────────────────────────────────────────────────────

from dataclasses import dataclass, field


@dataclass
class ExtractedUnit:
    locator: str   # e.g. "Page 3", "Slide 5", "Sheet: Revenue", "Row block 1-50"
    text: str
    kind: str = "text"  # text | table | slide | sheet | notes


@dataclass
class ExtractedDocument:
    filename: str
    file_type: str
    units: list = field(default_factory=list)
    error: str = ""

    def full_text(self) -> str:
        """Flatten all units into one string for LLM context building."""
        parts = []
        for u in self.units:
            parts.append(f"[{self.filename} — {u.locator}]\n{u.text}")
        return "\n\n".join(parts)

    def is_empty(self) -> bool:
        return not self.units and not self.error


# ── PDF ──────────────────────────────────────────────────────────────────

def extract_pdf(file_bytes: bytes, filename: str) -> ExtractedDocument:
    doc = ExtractedDocument(filename=filename, file_type="pdf")
    try:
        import pdfplumber
        import io
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    doc.units.append(ExtractedUnit(
                        locator=f"Page {i}", text=text, kind="text"))
                tables = page.extract_tables() or []
                for t_idx, table in enumerate(tables, start=1):
                    table_text = _table_to_text(table)
                    if table_text:
                        doc.units.append(ExtractedUnit(
                            locator=f"Page {i}, Table {t_idx}",
                            text=table_text, kind="table"))
    except Exception as e:
        # Fallback to PyMuPDF if pdfplumber fails (e.g. malformed PDF)
        try:
            import fitz  # pymupdf
            import io
            pdf = fitz.open(stream=file_bytes, filetype="pdf")
            for i, page in enumerate(pdf, start=1):
                text = page.get_text().strip()
                if text:
                    doc.units.append(ExtractedUnit(
                        locator=f"Page {i}", text=text, kind="text"))
            pdf.close()
        except Exception as e2:
            doc.error = f"Could not read PDF ({e}); fallback also failed ({e2})."
    return doc


def _table_to_text(table: list) -> str:
    rows = []
    for row in table:
        cells = [str(c).strip() if c is not None else "" for c in row]
        rows.append(" | ".join(cells))
    return "\n".join(r for r in rows if r.strip(" |"))


# ── DOCX ─────────────────────────────────────────────────────────────────

def extract_docx(file_bytes: bytes, filename: str) -> ExtractedDocument:
    doc = ExtractedDocument(filename=filename, file_type="docx")
    try:
        import docx
        import io
        d = docx.Document(io.BytesIO(file_bytes))

        section_num = 0
        current_heading = "Document start"
        buffer = []

        def flush():
            nonlocal buffer
            if buffer:
                doc.units.append(ExtractedUnit(
                    locator=f"Section: {current_heading}",
                    text="\n".join(buffer), kind="text"))
                buffer = []

        for para in d.paragraphs:
            style = (para.style.name or "").lower()
            text = para.text.strip()
            if not text:
                continue
            if "heading" in style:
                flush()
                section_num += 1
                current_heading = text
            else:
                buffer.append(text)
        flush()

        for t_idx, table in enumerate(d.tables, start=1):
            rows = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append(" | ".join(cells))
            table_text = "\n".join(r for r in rows if r.strip(" |"))
            if table_text:
                doc.units.append(ExtractedUnit(
                    locator=f"Table {t_idx}", text=table_text, kind="table"))

    except Exception as e:
        doc.error = f"Could not read DOCX ({e})."
    return doc


# ── PPTX ─────────────────────────────────────────────────────────────────

def extract_pptx(file_bytes: bytes, filename: str) -> ExtractedDocument:
    doc = ExtractedDocument(filename=filename, file_type="pptx")
    try:
        from pptx import Presentation
        import io
        prs = Presentation(io.BytesIO(file_bytes))

        for i, slide in enumerate(prs.slides, start=1):
            title = ""
            body_lines = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                is_title = (shape == slide.shapes.title) if slide.shapes.title else False
                if is_title and not title:
                    title = text
                else:
                    body_lines.append(text)

            slide_text = (f"Title: {title}\n" if title else "") + "\n".join(body_lines)
            if slide_text.strip():
                doc.units.append(ExtractedUnit(
                    locator=f"Slide {i}", text=slide_text.strip(), kind="slide"))

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes_text:
                    doc.units.append(ExtractedUnit(
                        locator=f"Slide {i} Notes", text=notes_text, kind="notes"))

    except Exception as e:
        doc.error = f"Could not read PPTX ({e})."
    return doc


# ── XLSX ─────────────────────────────────────────────────────────────────

def extract_xlsx(file_bytes: bytes, filename: str, max_rows_per_sheet: int = 200) -> ExtractedDocument:
    doc = ExtractedDocument(filename=filename, file_type="xlsx")
    try:
        import openpyxl
        import io
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count >= max_rows_per_sheet:
                    rows_text.append(f"... ({ws.max_row - max_rows_per_sheet} more rows truncated)")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows_text.append(" | ".join(cells))
                row_count += 1

            sheet_text = "\n".join(rows_text)
            if sheet_text.strip():
                doc.units.append(ExtractedUnit(
                    locator=f"Sheet: {sheet_name}", text=sheet_text, kind="sheet"))

    except Exception as e:
        doc.error = f"Could not read XLSX ({e})."
    return doc


# ── CSV ──────────────────────────────────────────────────────────────────

def extract_csv(file_bytes: bytes, filename: str, max_rows: int = 300) -> ExtractedDocument:
    doc = ExtractedDocument(filename=filename, file_type="csv")
    try:
        import pandas as pd
        import io
        df = pd.read_csv(io.BytesIO(file_bytes))
        truncated = len(df) > max_rows
        preview = df.head(max_rows)
        text = preview.to_csv(index=False)
        if truncated:
            text += f"\n... ({len(df) - max_rows} more rows truncated)"
        doc.units.append(ExtractedUnit(
            locator=f"Rows 1-{min(max_rows, len(df))}", text=text, kind="table"))
    except Exception as e:
        doc.error = f"Could not read CSV ({e})."
    return doc


# ── TXT (incl. transcripts / meeting notes) ───────────────────────────────

def extract_txt(file_bytes: bytes, filename: str) -> ExtractedDocument:
    doc = ExtractedDocument(filename=filename, file_type="txt")
    try:
        text = file_bytes.decode("utf-8", errors="replace").strip()
        if text:
            doc.units.append(ExtractedUnit(locator="Full text", text=text, kind="text"))
    except Exception as e:
        doc.error = f"Could not read text file ({e})."
    return doc


# ── Dispatcher ───────────────────────────────────────────────────────────

EXTENSION_MAP = {
    "pdf":  extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "xlsx": extract_xlsx,
    "csv":  extract_csv,
    "txt":  extract_txt,
}

SUPPORTED_EXTENSIONS = list(EXTENSION_MAP.keys())


def extract_file(uploaded_file) -> ExtractedDocument:
    """
    Accepts a Streamlit UploadedFile and routes it to the right extractor
    based on its extension. Returns an ExtractedDocument either way (with
    `.error` set if extraction failed), so callers never need a try/except.
    """
    filename = uploaded_file.name
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext not in EXTENSION_MAP:
        d = ExtractedDocument(filename=filename, file_type=ext or "unknown")
        d.error = f"Unsupported file type '.{ext}'. Supported: {', '.join(SUPPORTED_EXTENSIONS)}."
        return d

    file_bytes = uploaded_file.getvalue()
    extractor = EXTENSION_MAP[ext]
    return extractor(file_bytes, filename)
